"""
PPT parser service.
Parses PowerPoint (.pptx) files and generates training scripts from slide content.
"""
import logging
from typing import Optional, List, Dict, Any

from pptx import Presentation

logger = logging.getLogger(__name__)


def parse_ppt(file_path: str) -> Dict[str, Any]:
    """
    Parse a PowerPoint (.pptx) file and extract slide content.

    Args:
        file_path: Absolute path to the .pptx file.

    Returns:
        Dict with structure:
        {
            "slides": [{"title": "...", "content": "...", "notes": "..."}],
            "total_slides": N,
        }
    """
    prs = Presentation(file_path)
    slides = []

    for slide in prs.slides:
        title = ""
        content_parts: List[str] = []
        notes = ""

        # Extract title
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()

        # Extract body text from all shapes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Skip the title shape to avoid duplication
            if slide.shapes.title and shape.shape_id == slide.shapes.title.shape_id:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # Detect bullet level for formatting
                    level = paragraph.level or 0
                    prefix = "  " * level + "- " if level > 0 else ""
                    content_parts.append(prefix + text)

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "title": title,
            "content": "\n".join(content_parts),
            "notes": notes,
        })

    logger.info(f"Parsed PPT: {len(slides)} slides from {file_path}")

    return {
        "slides": slides,
        "total_slides": len(slides),
    }


def generate_script_from_ppt(
    ppt_data: Dict[str, Any],
    language: str = "hinglish",
) -> str:
    """
    Generate a training video script from parsed PPT slide data.

    Each slide becomes a section in the script with Hinglish-style narration
    transitions between slides.

    Args:
        ppt_data: Output of parse_ppt().
        language: Target language (hinglish, hindi, english).

    Returns:
        Generated script text.
    """
    slides = ppt_data.get("slides", [])
    total = ppt_data.get("total_slides", len(slides))

    if not slides:
        return "Namaste GroMo Partners! Is presentation mein koi content nahi mila."

    sections: List[str] = []

    for idx, slide in enumerate(slides):
        title = slide.get("title", "").strip()
        content = slide.get("content", "").strip()
        notes = slide.get("notes", "").strip()

        # Build section header
        slide_label = title if title else f"Slide {idx + 1}"

        # Add transition narration based on position
        if idx == 0:
            transition = _intro_transition(slide_label, language)
        elif idx == total - 1:
            transition = _outro_transition(language)
        else:
            transition = _mid_transition(idx, slide_label, language)

        # Build slide section
        section_parts = [f"[SLIDE {idx + 1}: {slide_label}]"]
        section_parts.append(transition)

        # Use speaker notes as narration if available, otherwise use content
        if notes:
            section_parts.append(notes)
        elif content:
            section_parts.append(content)
        elif title:
            section_parts.append(title)

        sections.append("\n".join(section_parts))

    return "\n\n".join(sections)


def _intro_transition(title: str, language: str) -> str:
    """Generate intro transition text."""
    if language == "english":
        return f"Hello GroMo Partners! Today we will learn about: {title}. Let's get started."
    if language == "hindi":
        return f"Namaste GroMo Partners! Aaj hum seekhenge: {title}. Chaliye shuru karte hain."
    # hinglish (default)
    return f"Namaste GroMo Partners! Aaj ka topic hai: {title}. Let's get started!"


def _mid_transition(index: int, title: str, language: str) -> str:
    """Generate mid-presentation transition text."""
    transitions_hinglish = [
        f"Ab baat karte hain: {title}.",
        f"Chaliye next point dekhte hain: {title}.",
        f"Aage badhte hain - ab dekhenge: {title}.",
        f"Ek important point hai: {title}.",
    ]
    transitions_english = [
        f"Now let's talk about: {title}.",
        f"Moving on to: {title}.",
        f"Next, let's look at: {title}.",
        f"An important point: {title}.",
    ]
    transitions_hindi = [
        f"Ab baat karte hain: {title}.",
        f"Aage badhte hain: {title}.",
        f"Ab dekhte hain: {title}.",
        f"Ek aur important point: {title}.",
    ]

    if language == "english":
        choices = transitions_english
    elif language == "hindi":
        choices = transitions_hindi
    else:
        choices = transitions_hinglish

    return choices[index % len(choices)]


def _outro_transition(language: str) -> str:
    """Generate outro transition text."""
    if language == "english":
        return (
            "So partners, that covers our presentation. "
            "Use this knowledge to help your customers better. Happy Selling!"
        )
    if language == "hindi":
        return (
            "Toh partners, yeh thi humari presentation. "
            "Is knowledge ka use karke apne customers ko best service dein. Happy Selling!"
        )
    # hinglish
    return (
        "Toh partners, yeh thi humari presentation. "
        "Isko use karke apne customers ko best service dein aur earnings badhayein! Happy Selling!"
    )
