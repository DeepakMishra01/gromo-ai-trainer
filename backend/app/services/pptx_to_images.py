"""
PPTX to Images converter.
Uses LibreOffice headless to convert PPTX → PDF, then PyMuPDF to render
each PDF page as a high-quality slide image.
This preserves ALL Gamma AI visuals — images, backgrounds, icons, formatting.
"""
import os
import shutil
import logging
import subprocess
import tempfile
from typing import Dict, List, Optional

from PIL import Image

logger = logging.getLogger(__name__)

SLIDE_W, SLIDE_H = 1280, 720
SOFFICE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"


def pptx_to_images(
    pptx_path: str,
    output_dir: str,
    width: int = SLIDE_W,
    height: int = SLIDE_H,
) -> List[str]:
    """
    Convert a PPTX file to a list of high-quality slide images.

    Pipeline: PPTX → PDF (LibreOffice) → per-page PNG (PyMuPDF) → resize

    Args:
        pptx_path: Path to the PPTX file.
        output_dir: Directory to save the output images.
        width: Target image width.
        height: Target image height.

    Returns:
        List of paths to the generated slide images (in slide order).
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    os.makedirs(output_dir, exist_ok=True)

    # Step 1: Convert PPTX → PDF using LibreOffice headless
    pdf_path = _convert_pptx_to_pdf(pptx_path)

    if not pdf_path or not os.path.exists(pdf_path):
        logger.warning("LibreOffice PDF conversion failed, falling back to qlmanage")
        return _fallback_qlmanage(pptx_path, output_dir, width, height)

    # Step 2: Render each PDF page as a high-quality PNG
    try:
        image_paths = _render_pdf_pages(pdf_path, output_dir, width, height)
    finally:
        # Cleanup temp PDF
        try:
            os.remove(pdf_path)
        except OSError:
            pass

    if not image_paths:
        logger.warning("PDF page rendering produced no images, using fallback")
        return _fallback_qlmanage(pptx_path, output_dir, width, height)

    # Step 3: Overlay GroMo branding on each slide
    image_paths = _add_gromo_branding(image_paths)

    logger.info(f"Converted {len(image_paths)} slides to images in {output_dir}")
    return image_paths


def _convert_pptx_to_pdf(pptx_path: str) -> Optional[str]:
    """
    Convert PPTX to PDF using LibreOffice headless mode.
    Returns path to the generated PDF, or None on failure.
    """
    if not os.path.exists(SOFFICE_PATH):
        logger.warning(f"LibreOffice not found at {SOFFICE_PATH}")
        return None

    tmp_dir = tempfile.mkdtemp(prefix="pptx_pdf_")

    try:
        result = subprocess.run(
            [
                SOFFICE_PATH,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", tmp_dir,
                pptx_path,
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if result.returncode != 0:
            logger.error(f"LibreOffice conversion failed: {result.stderr}")
            return None

        # Find the generated PDF
        basename = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(tmp_dir, f"{basename}.pdf")

        if os.path.exists(pdf_path):
            logger.info(f"PPTX → PDF: {os.path.getsize(pdf_path) / 1024:.0f}KB")
            return pdf_path

        # Check for any PDF in the output dir
        for f in os.listdir(tmp_dir):
            if f.endswith(".pdf"):
                return os.path.join(tmp_dir, f)

        return None

    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timed out")
        return None
    except Exception as e:
        logger.error(f"LibreOffice conversion error: {e}")
        return None


def _render_pdf_pages(
    pdf_path: str,
    output_dir: str,
    width: int = SLIDE_W,
    height: int = SLIDE_H,
) -> List[str]:
    """
    Render each page of a PDF as a high-quality PNG image using PyMuPDF.
    """
    import fitz  # PyMuPDF

    doc = fitz.open(pdf_path)
    num_pages = doc.page_count
    image_paths = []

    logger.info(f"Rendering {num_pages} PDF pages to images...")

    for i in range(num_pages):
        page = doc[i]
        out_path = os.path.join(output_dir, f"gamma_slide_{i + 1:03d}.png")

        try:
            # Render at 2.5x zoom for crisp quality
            zoom = 2.5
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)

            # Save high-res version
            pix.save(out_path)

            # Resize to target dimensions
            img = Image.open(out_path).convert("RGB")
            img = img.resize((width, height), Image.LANCZOS)
            img.save(out_path, "PNG", quality=95)

            image_paths.append(out_path)
            size_kb = os.path.getsize(out_path) / 1024
            logger.info(f"Slide {i + 1}: {width}x{height} ({size_kb:.0f}KB)")

        except Exception as e:
            logger.error(f"Failed to render PDF page {i + 1}: {e}")

    doc.close()
    return image_paths


def _fallback_qlmanage(
    pptx_path: str,
    output_dir: str,
    width: int,
    height: int,
) -> List[str]:
    """
    Fallback: use qlmanage to generate at least a thumbnail of the first slide.
    Then use python-pptx text extraction for remaining slides.
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    num_slides = len(prs.slides)
    image_paths = []

    # Try qlmanage for first slide thumbnail
    thumb_dir = tempfile.mkdtemp(prefix="pptx_thumbs_")
    try:
        result = subprocess.run(
            ["qlmanage", "-t", "-s", str(max(width, 1920)), "-o", thumb_dir, pptx_path],
            capture_output=True, text=True, timeout=30,
        )

        if result.returncode == 0:
            basename = os.path.basename(pptx_path)
            thumb = os.path.join(thumb_dir, f"{basename}.png")
            if os.path.exists(thumb):
                out_path = os.path.join(output_dir, "gamma_slide_001.png")
                img = Image.open(thumb).convert("RGB").resize((width, height), Image.LANCZOS)
                img.save(out_path, "PNG", quality=95)
                image_paths.append(out_path)
    except Exception:
        pass
    finally:
        shutil.rmtree(thumb_dir, ignore_errors=True)

    # For remaining slides, render text-based fallbacks
    from PIL import ImageDraw, ImageFont

    for i, slide in enumerate(prs.slides):
        if i == 0 and image_paths:
            continue  # Already have first slide from qlmanage

        out_path = os.path.join(output_dir, f"gamma_slide_{i + 1:03d}.png")
        img = _render_text_slide(slide, width, height, i + 1)
        img.save(out_path, "PNG", quality=95)
        image_paths.append(out_path)

    return image_paths


def _render_text_slide(slide, width: int, height: int, slide_num: int) -> Image.Image:
    """Render a slide as text on a gradient background."""
    from PIL import ImageDraw, ImageFont

    img = Image.new("RGB", (width, height))
    draw = ImageDraw.Draw(img)

    # Gradient background
    for y in range(height):
        r = int(15 + (25 - 15) * y / height)
        g = int(23 + (45 - 23) * y / height)
        b = int(42 + (90 - 42) * y / height)
        draw.line([(0, y), (width, y)], fill=(r, g, b))

    try:
        title_font = ImageFont.truetype("/System/Library/Fonts/HelveticaNeue.ttc", 36, index=1)
        body_font = ImageFont.truetype("/System/Library/Fonts/HelveticaNeue.ttc", 22)
    except Exception:
        title_font = ImageFont.load_default()
        body_font = ImageFont.load_default()

    title_text = ""
    content_lines = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    if not title_text:
                        title_text = text
                    else:
                        content_lines.append(text)

    if not title_text:
        title_text = f"Slide {slide_num}"

    draw.text((60, 80), title_text[:80], fill=(255, 255, 255), font=title_font)

    y = 160
    for line in content_lines[:12]:
        if y > height - 60:
            break
        draw.text((60, y), line[:100], fill=(200, 210, 220), font=body_font)
        y += 32

    return img


# ── GroMo Branding ──

_ASSETS_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
    "assets", "branding",
)


def _add_gromo_branding(image_paths: List[str]) -> List[str]:
    """
    Overlay GroMo logo and 'Powered by GroMo' badge on each slide image.
    Places the logo at top-left and badge at bottom-right.
    """
    logo_path = os.path.join(_ASSETS_DIR, "gromo_watermark.png")
    badge_path = os.path.join(_ASSETS_DIR, "powered_by_gromo.png")

    logo_img = None
    badge_img = None

    if os.path.exists(logo_path):
        try:
            logo_img = Image.open(logo_path).convert("RGBA")
            # Scale logo to fit nicely (max 140px wide)
            ratio = 140 / logo_img.width
            logo_img = logo_img.resize(
                (int(logo_img.width * ratio), int(logo_img.height * ratio)),
                Image.LANCZOS,
            )
        except Exception as e:
            logger.warning(f"Failed to load GroMo logo: {e}")

    if os.path.exists(badge_path):
        try:
            badge_img = Image.open(badge_path).convert("RGBA")
            # Scale badge
            ratio = 200 / badge_img.width
            badge_img = badge_img.resize(
                (int(badge_img.width * ratio), int(badge_img.height * ratio)),
                Image.LANCZOS,
            )
        except Exception as e:
            logger.warning(f"Failed to load GroMo badge: {e}")

    if not logo_img and not badge_img:
        return image_paths

    for img_path in image_paths:
        try:
            slide = Image.open(img_path).convert("RGBA")

            if logo_img:
                # Top-left with padding
                slide.paste(logo_img, (20, 12), logo_img)

            if badge_img:
                # Bottom-right with padding
                bx = slide.width - badge_img.width - 20
                by = slide.height - badge_img.height - 12
                slide.paste(badge_img, (bx, by), badge_img)

            slide.convert("RGB").save(img_path, "PNG", quality=95)
        except Exception as e:
            logger.warning(f"Failed to add branding to {img_path}: {e}")

    logger.info(f"Added GroMo branding to {len(image_paths)} slides")
    return image_paths


def extract_slide_texts(pptx_path: str) -> List[Dict]:
    """
    Extract title and content text from each slide in a PPTX.
    Used to generate per-slide narration that matches slide content.

    Returns list of dicts: [{"title": "...", "content": "...", "slide_num": N}, ...]
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    slides_text = []

    for i, slide in enumerate(prs.slides):
        title = ""
        content_parts = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # First substantial text is title
                if not title and len(text) > 3:
                    title = text
                else:
                    content_parts.append(text)

        slides_text.append({
            "slide_num": i + 1,
            "title": title,
            "content": " | ".join(content_parts[:10]),  # Join content points
        })

    return slides_text
